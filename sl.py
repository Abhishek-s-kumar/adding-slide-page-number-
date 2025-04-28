from pptx import Presentation
from pptx.util import Pt

# Load the presentation
prs = Presentation('PPT 11.pptx')

# Get the slide width and height
slide_width = prs.slide_width
slide_height = prs.slide_height

for idx, slide in enumerate(prs.slides):
    # Set margins (you can tweak)
    margin_right = Pt(50)  # 50 points from right
    margin_bottom = Pt(30) # 30 points from bottom

    # Calculate position
    left = slide_width - margin_right
    top = slide_height - margin_bottom
    width = Pt(50)
    height = Pt(20)

    # Add textbox for slide number
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(idx + 1)  # Slide number

    # Style
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.name = 'Arial'

    # Align to right
    p.alignment = 2  # 2 = RIGHT

# Save the modified presentation
prs.save('presentation_with_slide_numbers.pptx')

print("Slide numbers added manually at the bottom-right and saved as 'presentation_with_slide_numbers.pptx'.")
